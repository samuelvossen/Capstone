import os
from io import open
import shutil
import collections
import collections.abc
from os.path import basename
import pptx.util
from pptx import Presentation
from pptx.util import Pt
from bs4 import BeautifulSoup
import requests
import deepl

""" Text translation using DeepL API!!!
-To use the DeepL Python Library, you'll need an API authentication key.
To get a key, please create an account here --> https://www.deepl.com/docs-api/api-access/
With a DeepL API Free account you can translate up to 500,000 characters/month for free.

Installation
The library can be installed from PyPI using pip:

pip install --upgrade deepl """

# slideOption controls layout of powerpoint. This program supports options 1, 4, and 5. 
# 1 for default layout, 4 for landscape with text on the left, 5 for landscape with picture on the left

methods = "Web"
slideOption = 4
textFont = titleFont = "NATS"
textSize = 13
titleSize = 2
output_file = 'project_abcd.pptx'
dest_language = "ZH" #Chinese language

translation_engine = 'DeepL'

""" DeepL API integration"""
auth_key = "b756986e-67d1-762d-cb6f-0971053d969d:fx"  #(auth_key) Replace with your key
translator = deepl.Translator(auth_key)

#DeepL API testing
result = translator.translate_text("Hello, world!", target_lang="FR")
print(result.text)  # "Bonjour, le monde !"

do_translation = True

#all_pages = [26, 27, 28, 29, 30, 39, 50, 52, 53, 110, 111, 112, 116, 196, 206, 262, 265, 275, 276, 314, 317, 318, 319, 320, 321, 322, 324, 325, 326, 327, 328, 329, 401, 405, 406, 407, 409, 410, 411, 412, 413, 415, 418, 419, 422, 423, 424, 425, 426, 427, 428, 431, 432, 433, 434, 437, 438, 439, 440, 441, 442, 443, 444, 445, 462, 463, 468, 469, 470, 471, 472, 475, 476, 477, 478, 483, 484, 491, 492, 493, 502, 506, 520, 542, 544, 549, 568, 574, 578, 581, 582, 601, 605, 611, 626, 627, 631, 649, 654, 655, 658, 659, 660, 662, 664, 665, 666, 667, 670, 671, 672, 674, 678, 681, 682, 683, 684, 686, 688, 689, 690, 691, 693, 695, 696, 697, 698, 699, 700]

all_pages = [26, 27, 28, 29, 30, 39, 50, 52, 53, 110]

def buildPresentation():
    prs = Presentation()
    presentationLength = len(all_pages)
    pictureSlide = 0

    if (slideOption == 2):
        pictureSlide = 1
    # web scrapes the URL to get all needed information from the page
    for i in range(0, presentationLength):
        URL = "https://projectabcd.com/display_the_dress.php?id=" + str(all_pages[i])
        page = requests.get(URL, headers={"User-Agent": "html"})
        soup = BeautifulSoup(page.content, "html.parser")
        logo = soup.find("img")
        printLogo = logo.attrs["src"]
        logoURL = "http://projectabcd.com/" + printLogo
        r = requests.get(logoURL, headers={"User-Agent": "html"}, stream=True)
        if r.status_code == 200:
            with open(basename(printLogo), "wb") as f:
                r.raw.decode_content = True
                shutil.copyfileobj(r.raw, f)
        pageInfo = soup.find("div", class_="containerTitle")
        pageInfoImg = soup.find("div", class_="container")
        name = pageInfo.find("h2", class_="headTwo")
        printName = name.text
        image = pageInfoImg.find("div", class_="containerImage")
        img = image.find("image", class_="image")
        printImage = img.get("src")
        pictureURL = "http://projectabcd.com/" + printImage
        r = requests.get(pictureURL, headers={"User-Agent": "html"}, stream=True)
        if r.status_code == 200:
            with open(basename(printImage), "wb") as f:
                r.raw.decode_content = True
                shutil.copyfileobj(r.raw, f)

        pagetext = pageInfoImg.find("div", class_="containerText")
        description = pagetext.find("p", class_="words")
        printDescription = (description.text)
        fact = description.find_next_sibling("p")
        printFact = fact.text

        # We now got printDescription and printFact
        if do_translation:
             printDescription = DeepL(printDescription)
             printFact = DeepL(printFact)

        # creates the slide presentation if slide option 1 is choosen
        if (slideOption == 1):
            # creates the slides and sets layout preferences
            slide_layout = prs.slide_layouts[6]
            slide = prs.slides.add_slide(slide_layout)
            prs.slide_width = pptx.util.Inches(8)
            prs.slide_height = pptx.util.Inches(11)
            # places the logo on the slide
            logoHolder = slide.shapes.add_picture(basename(printLogo), pptx.util.Inches(7), pptx.util.Inches(0),
                                                  width=pptx.util.Inches(1), height=pptx.util.Inches(1))
            # places the title on the slide
            titleBox = slide.shapes.add_textbox(pptx.util.Inches(2.5), pptx.util.Inches(.5),
                                                width=pptx.util.Inches(3), height=pptx.util.Inches(1))
            titleBoxtf = titleBox.text_frame
            title = titleBoxtf.add_paragraph()
            title.text = printName
            title.font.name = titleFont
            title.font.size = Pt(titleSize)
            # places the picture on the slide
            pictureHolder = prs.slides[i].shapes
            pictureHolder.add_picture(basename(printImage), pptx.util.Inches(2.5), pptx.util.Inches(2),
                                      width=pptx.util.Inches(3), height=pptx.util.Inches(4))
            # creates a textbox for the description and fun fact
            contentBox = slide.shapes.add_textbox(pptx.util.Inches(1), pptx.util.Inches(6),
                                                  width=pptx.util.Inches(6), height=pptx.util.Inches(5))
            contentBoxtf = contentBox.text_frame
            contentBoxtf.word_wrap = True
            descriptionTitle = contentBoxtf.add_paragraph()
            descriptionTitle.font.name = textFont
            descriptionTitle.font.bold = True
            descriptionTitle.font.size = Pt(textSize)
            descriptionTitle.text = "Description: "
            descriptionParagraph = contentBoxtf.add_paragraph()
            descriptionParagraph.font.name = textFont
            descriptionParagraph.font.size = Pt(textSize)
            descriptionParagraph.text = printDescription
            FunFactTitle = contentBoxtf.add_paragraph()
            FunFactTitle.font.name = textFont
            FunFactTitle.font.bold = True
            FunFactTitle.font.size = Pt(textSize)
            FunFactTitle.text = "\nFun Fact:"
            FunFactParagraph = contentBoxtf.add_paragraph()
            FunFactParagraph.font.name = textFont
            FunFactParagraph.font.size = Pt(textSize)
            FunFactParagraph.text = printFact
        # creates the slide presentation if slide option 4 is chosen
        elif (slideOption == 4):
             # creates slide preferences
            slide_layout = prs.slide_layouts[6]
            prs.slide_width = pptx.util.Inches(11)
            prs.slide_height = pptx.util.Inches(8)
            slide2 = prs.slides.add_slide(slide_layout)
            # places picture to cover whole slide
            pictureHolder = prs.slides[pictureSlide].shapes
            pictureHolder.add_picture(basename(printImage), pptx.util.Inches(6), pptx.util.Inches(1),
                                      width=pptx.util.Inches(4), height=pptx.util.Inches(6))
            # creates next slide

            # place logo on the slide
            logoHolder = slide2.shapes.add_picture(basename(printLogo), pptx.util.Inches(5.5), pptx.util.Inches(0),
                                                   width=pptx.util.Inches(1), height=pptx.util.Inches(1))
            # places the title
            titleBox = slide2.shapes.add_textbox(pptx.util.Inches(4), pptx.util.Inches(1.5),
                                                 width=pptx.util.Inches(2), height=pptx.util.Inches(1))
            titleBoxtf = titleBox.text_frame
            title = titleBoxtf.add_paragraph()
            title.text = printName
            title.font.size = Pt(titleSize)
            title.font.name = titleFont
            # creates textbox for description and fun fact
            contentBox = slide2.shapes.add_textbox(pptx.util.Inches(1), pptx.util.Inches(1),
                                                   width=pptx.util.Inches(3), height=pptx.util.Inches(4))
            contentBoxtf = contentBox.text_frame
            contentBoxtf.word_wrap = True
            descriptionTitle = contentBoxtf.add_paragraph()
            descriptionTitle.font.name = textFont
            descriptionTitle.font.bold = True
            descriptionTitle.font.size = Pt(textSize)
            descriptionTitle.text = "Description: "
            descriptionParagraph = contentBoxtf.add_paragraph()
            descriptionParagraph.font.name = textFont
            descriptionParagraph.font.size = Pt(textSize)
            descriptionParagraph.text = printDescription
            FunFactTitle = contentBoxtf.add_paragraph()
            FunFactTitle.font.bold = True
            FunFactTitle.font.name = textFont
            FunFactTitle.font.size = Pt(textSize)
            FunFactTitle.text = "\nFun Fact:"
            FunFactParagraph = contentBoxtf.add_paragraph()
            FunFactParagraph.font.name = textFont
            FunFactParagraph.font.size = Pt(textSize)
            FunFactParagraph.text = printFact
            pictureSlide = pictureSlide + 1

        # creates the slide presentation if slide option 5 is chosen
        elif (slideOption == 5):
            # creates slide preferences
            slide_layout = prs.slide_layouts[6]
            prs.slide_width = pptx.util.Inches(11)
            prs.slide_height = pptx.util.Inches(8)
            slide2 = prs.slides.add_slide(slide_layout)
            # places picture to cover whole slide
            pictureHolder = prs.slides[pictureSlide].shapes
            pictureHolder.add_picture(basename(printImage), pptx.util.Inches(1), pptx.util.Inches(1),
                                      width=pptx.util.Inches(4), height=pptx.util.Inches(6))
            # creates next slide

            # place logo on the slide
            logoHolder = slide2.shapes.add_picture(basename(printLogo), pptx.util.Inches(5), pptx.util.Inches(0),
                                                   width=pptx.util.Inches(1), height=pptx.util.Inches(1))
            # places the title
            titleBox = slide2.shapes.add_textbox(pptx.util.Inches(4), pptx.util.Inches(1.5),
                                                 width=pptx.util.Inches(2), height=pptx.util.Inches(1))
            titleBoxtf = titleBox.text_frame
            title = titleBoxtf.add_paragraph()
            title.text = printName
            title.font.size = Pt(titleSize)
            title.font.name = titleFont
            # creates textbox for description and fun fact
            contentBox = slide2.shapes.add_textbox(pptx.util.Inches(7), pptx.util.Inches(1),
                                                   width=pptx.util.Inches(3), height=pptx.util.Inches(4))
            contentBoxtf = contentBox.text_frame
            contentBoxtf.word_wrap = True
            descriptionTitle = contentBoxtf.add_paragraph()
            descriptionTitle.font.name = textFont
            descriptionTitle.font.bold = True
            descriptionTitle.font.size = Pt(textSize)
            descriptionTitle.text = "Description: "
            descriptionParagraph = contentBoxtf.add_paragraph()
            descriptionParagraph.font.name = textFont
            descriptionParagraph.font.size = Pt(textSize)
            descriptionParagraph.text = printDescription
            FunFactTitle = contentBoxtf.add_paragraph()
            FunFactTitle.font.bold = True
            FunFactTitle.font.name = textFont
            FunFactTitle.font.size = Pt(textSize)
            FunFactTitle.text = "\nFun Fact:"
            FunFactParagraph = contentBoxtf.add_paragraph()
            FunFactParagraph.font.name = textFont
            FunFactParagraph.font.size = Pt(textSize)
            FunFactParagraph.text = printFact
            pictureSlide = pictureSlide + 1

    prs.save(output_file)
    return output_file

def DeepL(printDescription):
    result = (translator.translate_text(printDescription,target_lang=dest_language,split_sentences=1)).__str__().encode("utf-8")

    return result


ppt_file = buildPresentation()
os.startfile(ppt_file)
